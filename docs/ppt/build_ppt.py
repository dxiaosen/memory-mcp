"""Memory MCP 答辩 PPT —— 极简版。
每页正文内容都用 matplotlib 画成一张图插入，PPT 只放标题+图。
这彻底避免了原生形状手动堆坐标导致的错位/遮挡问题。
配色：主蓝 #005982、红 #D60012、灰 #666666。字体 Noto Sans CJK SC。
中文引述一律用「」避免破坏 Python 字符串边界。
"""
from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
IMG = ROOT / "docs" / "ppt" / "imgs"
BG = Path("/tmp/tpl_bg")
OUT = ROOT / "docs" / "ppt" / "memory-mcp-defense.pptx"

SW, SH = Inches(13.333), Inches(7.5)
NAVY = RGBColor(0x00, 0x59, 0x82)
RED = RGBColor(0xD6, 0x00, 0x12)
MID = RGBColor(0x66, 0x66, 0x66)
INK = RGBColor(0x26, 0x2B, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Noto Sans CJK SC"


def _fill_pic(slide, name):
    """插入 imgs/ 下图片到正文区，等比缩放居中，不超出正文区。"""
    from PIL import Image as PILImage
    im = PILImage.open(IMG / f"{name}.png")
    iw, ih = im.size
    max_w, max_h = 12.1, 5.6
    w = max_w
    h = w * ih / iw
    if h > max_h:  # 超高则按高度反算
        h = max_h; w = h * iw / ih
    x = (13.333 - w) / 2
    y = 1.05 + (5.75 - h) / 2
    if y < 1.0: y = 1.0
    slide.shapes.add_picture(str(IMG / f"{name}.png"), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def bg(slide, name):
    slide.shapes.add_picture(str(BG / f"{name}.png"), 0, 0, width=SW, height=SH)


def text(slide, val, x, y, w, h, size=18, color=INK, bold=False,
         align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    tf.margin_left = Inches(0.02); tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]; p.text = val; p.alignment = align
    p.space_after = Pt(0); p.space_before = Pt(0); p.line_spacing = 1.08
    for r in p.runs:
        r.font.name = FONT; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return box


def page_no(slide, label):
    text(slide, label, 11.85, 7.0, 0.9, 0.25, 8.5, MID, align=PP_ALIGN.RIGHT)


def notes(slide, secs, script):
    m, s = divmod(secs, 60)
    slide.notes_slide.notes_text_frame.text = f"【建议用时 {m:02d}:{s:02d}】\n{script.strip()}"


def content_slide(prs, title, img, label, secs=60, script=""):
    """标准正文页：背景 + 标题 + 分隔线 + 一张图 + 页码。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6]); bg(slide, "content_a")
    text(slide, title, 0.62, 0.34, 12.1, 0.55, 22, NAVY, True, valign=MSO_ANCHOR.MIDDLE)
    # 分隔线
    ln = slide.shapes.add_shape(1, Inches(0.62), Inches(0.92), Inches(12.1), Pt(2))
    ln.fill.solid(); ln.fill.fore_color.rgb = NAVY; ln.line.fill.background(); ln.shadow.inherit = False
    _fill_pic(slide, img)
    page_no(slide, label)
    notes(slide, secs, script)
    return slide


def section_slide(prs, num, title, subtitle, label, secs=8):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); bg(slide, "section")
    text(slide, num, 0.7, 1.0, 5.8, 3.2, 120, RED, True, valign=MSO_ANCHOR.MIDDLE)
    text(slide, title, 7.6, 2.5, 5.1, 0.9, 30, NAVY, True, valign=MSO_ANCHOR.MIDDLE)
    text(slide, subtitle, 7.65, 3.55, 5.0, 0.7, 15, RED)
    page_no(slide, label)
    notes(slide, secs, f"转场：下面进入「{title}」。")
    return slide


def build():
    prs = Presentation(); prs.slide_width = SW; prs.slide_height = SH
    p = prs.core_properties; p.title = "Memory MCP 答辩"; p.subject = "owner-scoped 长期记忆服务"

    # ===== 第一幕 痛点与背景 =====
    # P1 封面
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, "cover")
    text(s, "Memory MCP", 0.92, 2.28, 6.8, 0.7, 38, NAVY, True)
    text(s, "让 Agent 记住判断，也记住边界", 0.95, 2.98, 7.2, 0.45, 22, NAVY, True)
    text(s, "跨会话长期记忆服务 · 答辩汇报", 0.95, 3.92, 6.0, 0.32, 15, RED, True)
    text(s, "汇报人：[姓名]      部门：[部门]", 0.95, 4.32, 6.5, 0.3, 13, INK)
    text(s, "日期：[答辩日期]", 0.95, 4.66, 6.0, 0.3, 13, INK)
    notes(s, 25, "各位领导好，我汇报的课题是 Memory MCP。它解决 Agent 跨会话记忆三个问题：记不住、容易串、难治理。今天 20 分钟重点讲痛点与背景、和主流方案的差异化、核心机制、泡泡玛特双用户实测。")

    # P02 痛点
    content_slide(prs, "Agent 记忆的真正难点：不是存不存，而是存进去能不能管", "P02", "02", 65,
                  "长期记忆难点不在存储本身。第一是忘，重复建上下文损耗效率；第二是串，身份边界失守是信任事故且不可逆；第三是乱，错误结论无法修订作废，越积越脏。故目标是可治理的记忆基础设施，而非多存文本。")

    # P03 根因与对策：为什么独立成服务
    content_slide(prs, "痛点的根因：记忆散落在各 Agent 内，无统一治理", "P03", "03", 55,
                  "承接 P02 三个痛点，根因不在存储能力，而在记忆散落各 Agent、无统一治理。对策是将记忆抽成独立服务，由服务端统一负责抽取、准入、召回、身份隔离、生命周期，Agent 经 MCP 标准协议接入。类比 App 自建数据库 vs 统一 DB 服务——记忆复杂度在治理不在存储，此为立项依据。架构总览下一幕展开。")

    # ===== 第二幕 竞品对比与差异化 =====
    section_slide(prs, "02", "竞品对比与差异化", "判断演进审计 · 团队自动共识 · 失效治理", "04")

    # P05 竞品对比表（含最接近同行 TencentDB）
    content_slide(prs, "四款主流方案对比：即便最接近的 TencentDB，仍差在三轴", "P05", "05", 75,
                  "承接 P03「记忆该独立成服务」——市面上是否已有此类方案？看此表。ChatGPT 与 Mem0 仍为 SDK 或平台绑定、扁平事实、无团队。TencentDB Agent Memory 已是独立 MCP 服务、结构化、团队级——此为最接近的同行，予以承认。差异化收窄至站得住的三轴：①判断演进留审计链（TencentDB 仅资产版本号、无判断间 provenance）；②团队共识自动提取（其靠手动共享、Memory MCP 主动发现）；③失效治理（其无到期 TTL、Memory MCP 有准入+生命周期+脱敏）。勿逐格念，指这三轴结论。")

    # P06 三个差异化
    content_slide(prs, "三个核心差异化（对照 TencentDB）", "P06", "06", 60,
                  "三条差异化，每张对照 TencentDB 缺口与 Memory MCP 补位。第一判断演进留审计链——技术含量最高，一句话点透「改判断不改历史」；第二团队共识自动提取，主动发现而非手动共享；第三失效治理，到期 TTL+准入+生命周期。此三轴连最接近的 TencentDB 亦未覆盖。")

    # ===== 第三幕 核心机制 =====
    section_slide(prs, "03", "核心机制", "架构 · 闭环 · 数据模型 · 准入 · 召回 · 生命周期 · 团队 · 隔离", "07")

    # P08 系统架构总览
    content_slide(prs, "系统架构：Agent 只管接入，记忆治理全在服务端", "P08", "08", 55,
                  "架构总览：用户→三种宿主→Agent Client（BeforeRun 召回/AfterRun 捕获/多宿主适配/幂等去重）→Memory MCP Server（13 工具 + 准入/召回/生命周期/团队/身份隔离/脱敏审计）→PostgreSQL/LLM/Worker。核心设计：记忆归属由身份决定，与 Agent 实现解耦。下面逐层展开闭环。")

    # P09 核心闭环
    content_slide(prs, "核心闭环：写入异步治理，读取同步注入", "P09", "09", 65,
                  "闭环分写读。写入：对话完成→Stop Hook 强制每轮入队→Worker 异步结构化抽取候选→准入三级决策（auto_save/pending/discard）+ replacement 子动作→事务化落库。读取：新轮开始→BeforeRun Hook→SQL 内 owner 身份过滤→词法+向量+近期三路召回→关系感知补漏加权→Profile 优先级排序在 token 预算内注入。读写解耦保证用户不等抽取，身份过滤在召回前防越权。强调模型不需自判存不存，Hook 强制入队是可靠性根基。失败走 fail-open 不阻断 Agent，同 turn 重入靠 event_id 幂等兜底。")

    # P10 记忆数据模型
    content_slide(prs, "记忆数据模型：五张表，两层结构", "P10", "10", 60,
                  "这是数据模型本身。5 张表分两层：memory_items 持稳定身份（owner/profile/subject/memory_type 跨版本不变），memory_revisions 持版本快照（content/assertion_kind/valid 窗口/embedding 每次修订追加），两者 1:N——同一逻辑判断的多版本共存，旧版不删只标 superseded。周围挂三张表：memory_evidence 来源溯源（source_turn_id/source_expression），memory_reviews 待确认候选（status: pending/confirmed/rejected，确认后 resolved_memory_id 解析为 item），memory_relations 记忆间关系自引用（source/target→items，relation_type 如 supports/challenges）。底部标 memory_type（投研 8 类此处 4 代表）、lifecycle 四态、assertion_kind 三类。这不是扁平 key-value，是带立场、版本、来源的结构化判断。")

    # P11 准入四类
    content_slide(prs, "准入四类结果：不把什么都存，也不都让人审", "P11", "11", 55,
                  "核心卖点：不把什么都存也不都让人审。auto_save 自动沉淀高频偏好，pending 只打扰人一次。discard 过滤闲聊和模型产出。replacement 取代旧判断。提一句双重去重不展开。")

    # P12 召回三路
    content_slide(prs, "召回三路打分：向量 + 词法 + 关系，并行融合", "P12", "12", 60,
                  "召回不是只靠向量。三路并行：向量管语义近义、词法管精确命中与主题键、关系管判断间因果（supports/challenges）。任一路召回都进候选，加权融合后排序注入。身份过滤前置防越权。底部召回优先级 Profile 可调。一句话：RAG 只有一路向量，Memory MCP 是有治理的多路召回。")

    # P13 生命周期
    content_slide(prs, "生命周期：判断能改、能废、能复活", "P13", "13", 55,
                  "对应痛点第三条无治理。状态机：active 可被取代/撤销/到期。核心卖点是 revoke 后槽位释放可重建不死锁。replacement 保留旧版本可追溯判断怎么演进的。领导关心存错了怎么办——能改能废留痕能重建。")

    # P14 团队记忆
    content_slide(prs, "团队记忆：从个人共识到团队共识", "P14", "14", 60,
                  "团队记忆是竞品完全空白的能力。流程：周期性扫描成员个人记忆→embedding 聚类相似判断→团队待确认候选→任一成员确认→落 team owner 全员可见。幂等防重复，弱方向校验防把立场相反的判断并成一条。一句话：两人研究判断趋同时，系统主动发现并提议沉淀成团队共识。")

    # P15 身份隔离
    content_slide(prs, "身份隔离与安全：服务端强制，不靠客户端自觉", "P15", "15", 45,
                  "领导关心多用户会不会串。答案：owner 从 Token 派生，工具参数拒收 owner 防伪造。个人/团队多层隔离互不可写。安全三点：数据不丢、出错可查、安全合规。一页带过提问再展开。")

    # ===== 第四幕 泡泡玛特实测 =====
    section_slide(prs, "04", "泡泡玛特实测", "双用户 · 38 轮 · 真实对话验证", "16")

    # P17 测试设计
    content_slide(prs, "测试设计：两用户聊同一标的，殊途同归", "P17", "17", 60,
                  "测试设计：两用户聊同一标的泡泡玛特。用户一从出海切入，用户二从爆款依赖切入，都收敛到平台化验证判断——这让团队提取能聚类出共性候选。38 轮全经 Hook 自动捕获无人为干预。")

    # P18 演示截图
    content_slide(prs, "演示：一条判断的完整生命", "P18", "18", 70,
                  "全场高潮。三段真实对话：建立→auto_save 成 thesis；修订→replacement 取代旧版；团队共识→两用户收敛后 confirm 成团队记忆。用真实对话证明判断能演进、能跨人收敛、能沉淀。每段 1 分钟节奏要快。")

    # ===== 第五幕 持续演进收尾 =====
    content_slide(prs, "总结与持续演进方向", "P19", "19", 50,
                  "总结：做成了 MCP 标准、身份隔离、判断可演进、团队可共识的记忆服务。三个差异化：协议级、带立场的判断、判断演进+团队共识。后续：批量撤销、timeline 修复、更多 Profile。坦诚局限一句话带过。")

    # ===== 致谢 =====
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, "thanks")
    text(s, "致 谢", 0.72, 2.0, 4.0, 1.3, 52, NAVY, True)
    text(s, "为信任 责任 贡献 回报", 0.78, 3.5, 4.0, 0.4, 16, NAVY, True)
    text(s, "恳请各位领导批评指正", 0.78, 4.0, 5.0, 0.4, 18, RED, True)
    notes(s, 15, "汇报结束，恳请领导提问。")

    # ===== 备答页 Q&A =====
    content_slide(prs, "提问预判：高频问题备答", "PA1", "A1", 0,
                  "回答用业务风险和复用价值，不陷入实现细节。")

    return prs


def sanitize(path: Path, n_slides: int) -> None:
    """清理 python-pptx 默认模板残留，避免 PowerPoint 跨机器「需修复/打不开」。

    彻底处理：
    - 删 printerSettings、customXml 残留
    - 清 presentation.xml.rels 里的 printerSettings 引用
    - app.xml Slides 计数对齐
    - presProps 清 extLst
    - [Content_Types].xml 只保留实际存在的扩展名声明，删多余 Override
    """
    import re, shutil, zipfile
    tmp = path.with_suffix(".tmp.pptx")
    with zipfile.ZipFile(path) as zin, \
         zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        existing = set(zin.namelist())
        for name in existing:
            # 跳过打印机设置、自定义 XML 残留
            if "printerSettings" in name or name.startswith("customXml"):
                continue
            data = zin.read(name)
            if name == "ppt/_rels/presentation.xml.rels":
                txt = data.decode("utf-8")
                txt = re.sub(r'<Relationship[^>]*printerSettings[^>]*/>', "", txt)
                data = txt.encode("utf-8")
            elif name == "docProps/app.xml":
                txt = data.decode("utf-8")
                txt = re.sub(r"<Slides>\d+</Slides>", f"<Slides>{n_slides}</Slides>", txt)
                data = txt.encode("utf-8")
            elif name == "ppt/presProps.xml":
                txt = data.decode("utf-8")
                txt = re.sub(r"<p:extLst>.*?</p:extLst>", "", txt, flags=re.S)
                data = txt.encode("utf-8")
            elif name == "[Content_Types].xml":
                txt = data.decode("utf-8")
                # 删指向不存在文件的 Override（如 customXml、printerSettings）
                for m in re.finditer(r'<Override PartName="(/[^"]*)"[^>]*/>', txt):
                    part = m.group(1).lstrip("/")
                    if part not in existing and "customXml" not in part and "printerSettings" not in part:
                        # 只在确实不存在时才删（保守：保留 slideMaster/layout/notes）
                        if part.startswith("ppt/") and not any(part.startswith(p) for p in
                            ("ppt/slides/", "ppt/slideMasters", "ppt/slideLayouts", "ppt/notesMasters", "ppt/theme", "ppt/media", "ppt/embeddings")):
                            txt = txt.replace(m.group(), "")
                data = txt.encode("utf-8")
            zout.writestr(name, data)
    shutil.move(tmp, path)


if __name__ == "__main__":
    prs = build(); n = len(prs.slides); prs.save(OUT)
    sanitize(OUT, n)
    print(f"generated: {OUT}")
    print(f"slides: {n}")
